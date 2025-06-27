import os
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def parse_robot_output(xml_path):
    tree = ET.parse(xml_path)
    root = tree.getroot()

    total = int(root.attrib.get('total', 0))
    passed = int(root.attrib.get('passed', 0))
    failed = int(root.attrib.get('failed', 0))

    if total == 0 and root.find('statistics') is not None:
        stats = root.find('statistics')
        total = passed = failed = 0
        for stat in stats.findall('.//stat'):
            total += int(stat.attrib.get('total', 0))
            passed += int(stat.attrib.get('pass', 0))
            failed += int(stat.attrib.get('fail', 0))

    return total, passed, failed

def create_ppt(total, passed, failed):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # layout em branco

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Resumo dos Testes Robot Framework"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True

    p1 = content_tf.add_paragraph()
    p1.text = f"Total de testes: {total}"
    p1.font.size = Pt(24)
    p1.font.color.rgb = RGBColor(0, 0, 0)

    p2 = content_tf.add_paragraph()
    p2.text = f"Testes aprovados: {passed}"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0, 128, 0)

    p3 = content_tf.add_paragraph()
    p3.text = f"Testes falhados: {failed}"
    p3.font.size = Pt(24)
    p3.font.color.rgb = RGBColor(255, 0, 0)

    return prs

def add_screenshots_slide(prs, screenshots_folder):
    if not os.path.exists(screenshots_folder):
        print(f"Pasta de screenshots '{screenshots_folder}' não existe.")
        return

    imagens = [f for f in os.listdir(screenshots_folder) if f.lower().endswith('.png')]
    if not imagens:
        print("Nenhuma screenshot encontrada para adicionar ao PPT.")
        return

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.text = "Screenshots de Testes"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    top = Inches(1.5)
    left = Inches(1)
    max_height = Inches(3)

    for img_name in imagens:
        img_path = os.path.join(screenshots_folder, img_name)
        slide.shapes.add_picture(img_path, left, top, height=max_height)
        top += max_height + Inches(0.3)

if __name__ == "__main__":
    total, passed, failed = parse_robot_output('output.xml')
    prs = create_ppt(total, passed, failed)
    add_screenshots_slide(prs, 'screenshots')
    prs.save('test_report.pptx')
    print("PPT com resumo e screenshots salvo em test_report.pptx")
